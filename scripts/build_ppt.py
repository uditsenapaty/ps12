#!/usr/bin/env python
"""Fill the ISRO BAH 2026 idea-submission template with the PS-12 idea — world-class, professional.

Designed FOR the template: black header band (logos) + WHITE body + orange->purple->blue bottom rule.
Clean light slides, dark text, tasteful accents (ISRO orange / blue / violet), thin arrow-headed
connectors, a research-paper U-Net figure, and a realistic dashboard mock. No left-bar cards, no
"Thank You". Run:  python scripts/build_ppt.py  -> "ISRO_BAH_2026_PS12 - FILLED.pptx"
"""
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- design tokens ------------------------------------------------------------
FONT = "Segoe UI"
INK = RGBColor(0x16, 0x22, 0x3A); SLATE = RGBColor(0x55, 0x63, 0x77); MUTE = RGBColor(0x8A, 0x97, 0xA8)
HAIR = RGBColor(0xE3, 0xE9, 0xF0); MIST = RGBColor(0xF4, 0xF7, 0xFB); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF2, 0x6A, 0x1B); BLUE = RGBColor(0x27, 0x6E, 0xF0); VIOLET = RGBColor(0x6B, 0x4E, 0xE0)
TEAL = RGBColor(0x10, 0xA5, 0x9A); GREEN = RGBColor(0x1B, 0xA1, 0x66); NAVY = RGBColor(0x12, 0x1E, 0x36)
# encoder blue gradient (shallow->deep) / decoder teal gradient
ENC = [RGBColor(0x9D, 0xC2, 0xFF), RGBColor(0x5E, 0x97, 0xF5), RGBColor(0x35, 0x6F, 0xE0), RGBColor(0x1E, 0x49, 0xB8)]
DEC = [RGBColor(0x2C, 0x9B, 0x93), RGBColor(0x46, 0xB3, 0xA9), RGBColor(0x7C, 0xCF, 0xC7)]
IRBG = RGBColor(0x1A, 0x24, 0x36)

SRC = [s for s in glob.glob(r"D:\Udit\gitclones\ps12\*.pptx") if "FILLED" not in s][0]
prs = Presentation(SRC)
S = list(prs.slides)


# ---- primitives ---------------------------------------------------------------
def _ns(sh):
    sh.shadow.inherit = False
    return sh


def _run(p, t, s, c, b=False, i=False):
    r = p.add_run(); r.text = t
    f = r.font; f.size = Pt(s); f.color.rgb = c; f.bold = b; f.italic = i; f.name = FONT
    return r


def settext(sh, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tf = sh.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    for j, ln in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align); p.space_after = Pt(ln.get("sa", 3)); p.space_before = Pt(0)
        if isinstance(ln.get("runs"), list):
            for rr in ln["runs"]:
                _run(p, rr["t"], rr.get("s", 12), rr.get("c", INK), rr.get("b", False), rr.get("i", False))
        else:
            _run(p, ln["t"], ln.get("s", 12), ln.get("c", INK), ln.get("b", False), ln.get("i", False))
    return sh


def rect(sl, l, t, w, h, fill, line=None, lw=0.75, rounded=True, radius=0.08):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    if rounded:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return _ns(sh)


def oval(sl, l, t, w, h, fill, line=None):
    o = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    o.fill.solid(); o.fill.fore_color.rgb = fill
    if line is None:
        o.line.fill.background()
    else:
        o.line.color.rgb = line; o.line.width = Pt(1)
    return _ns(o)


def conn(sl, x1, y1, x2, y2, color=INK, w=1.3, head=True, dash=None):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    if head:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    return _ns(c)


def tb(sl, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    settext(box, lines, align=align, anchor=anchor)
    return box


def head(sl, kicker, title, accent=ORANGE):
    tb(sl, 0.52, 0.72, 9.0, 0.26, [{"t": kicker.upper(), "s": 10.5, "c": accent, "b": True}])
    tb(sl, 0.5, 0.98, 9.0, 0.56, [{"t": title, "s": 23, "c": INK, "b": True}])
    rect(sl, 0.54, 1.55, 1.25, 0.05, accent, rounded=False)


def badge(sl, cx, cy, txt, color, d=0.32):
    o = oval(sl, cx - d / 2, cy - d / 2, d, d, color)
    settext(o, [{"t": str(txt), "s": 11.5, "c": WHITE, "b": True, "align": PP_ALIGN.CENTER}],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def clear_prompts(sl):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.shape_type == 17:
            sh.text_frame.clear()


# ===============================================================================
# SLIDE 1 — cover fields (white bottom band -> DARK text)
for sh in S[0].shapes:
    if sh.has_text_frame and sh.text_frame.text.strip():
        tx = sh.text_frame.text
        if tx.startswith("Problem Statement"):
            settext(sh, [{"runs": [
                {"t": "Problem Statement:  ", "s": 13, "c": INK, "b": True},
                {"t": "PS-12 — Fill in the Frames Seamlessly: Enhancing Temporal Resolution of Satellite "
                      "Imagery using AI/ML based on Optical Flow", "s": 13, "c": NAVY}]}],
                anchor=MSO_ANCHOR.MIDDLE)
        elif tx.startswith("Team Name"):
            settext(sh, [{"runs": [{"t": "Team Name:  ", "s": 13, "c": INK, "b": True},
                                   {"t": "[ your team ]", "s": 13, "c": SLATE}]}], anchor=MSO_ANCHOR.MIDDLE)
        elif tx.startswith("Team Leader"):
            settext(sh, [{"runs": [{"t": "Team Leader:  ", "s": 13, "c": INK, "b": True},
                                   {"t": "[ your name ]", "s": 13, "c": SLATE}]}], anchor=MSO_ANCHOR.MIDDLE)

# ===============================================================================
# SLIDE 3 — Opportunity & USP  (two clean columns, no boxes)
sl = S[2]; clear_prompts(sl)
head(sl, "01 · Opportunity & Edge", "Minutes matter — we fill the gap")
tb(sl, 0.55, 1.85, 4.3, 3.3, [
    {"t": "THE GAP", "s": 11, "c": ORANGE, "b": True, "sa": 3},
    {"t": "Geostationary INSAT-3DR/3DS image every 30 minutes. Cyclones, thunderstorms, fire fronts "
          "and floods evolve in minutes — so the action happens between frames.", "s": 12, "c": SLATE, "sa": 10},
    {"t": "THE IDEA", "s": 11, "c": BLUE, "b": True, "sa": 3},
    {"t": "Teach an AI to synthesise the missing in-between frames from two real frames → 30 → 15 → "
          "7.5 min, with no new satellite hardware.", "s": 12, "c": SLATE, "sa": 10},
    {"t": "WHY WE'RE DIFFERENT", "s": 11, "c": VIOLET, "b": True, "sa": 3},
    {"t": "Classical optical flow (TV-L1) assumes straight, brightness-constant motion → it blurs and "
          "ghosts on fast, non-linear cloud growth. We learn the motion field and the fusion from the "
          "satellite thermal-IR itself.", "s": 12, "c": SLATE},
])
conn(sl, 5.05, 1.95, 5.05, 5.0, HAIR, w=1.0, head=False)
usp = [
    "Trained on satellite IR (brightness temperature) — not natural video.",
    "Cross-satellite transfer: learn on dense GOES-19 / Himawari (10-min) → apply to INSAT.",
    "Self-supervised INSAT adaptation — needs NO labels (uses INSAT's own 30-min frames).",
    "Custom UNetVFI = intermediate-flow (RIFE-style) + visibility blending (Super-SloMo-style).",
    "Complete product: .nc → .nc, web dashboard, and a metric-validated report.",
]
tb(sl, 5.3, 1.85, 4.25, 0.3, [{"t": "UNIQUE ADVANTAGES", "s": 11, "c": GREEN, "b": True}])
for i, u in enumerate(usp):
    y = 2.2 + i * 0.5
    oval(sl, 5.32, y + 0.05, 0.12, 0.12, GREEN)
    tb(sl, 5.55, y - 0.05, 4.0, 0.5, [{"t": u, "s": 11.5, "c": INK}])
rect(sl, 5.3, 4.78, 4.25, 0.42, MIST, line=HAIR, radius=0.18)
settext(sl.shapes[-1], [{"runs": [{"t": "Already real:  ", "s": 11, "c": GREEN, "b": True},
        {"t": "trained on a Tesla T4, validated on real GOES-19 (val PSNR ≈ 42, SSIM ≈ 0.95).",
         "s": 11, "c": INK}]}], anchor=MSO_ANCHOR.MIDDLE)

# ===============================================================================
# SLIDE 4 — Features (clean 3x2 tiles, numbered, no left bars)
sl = S[3]; clear_prompts(sl)
head(sl, "02 · Capabilities", "What the solution does", accent=BLUE)
feats = [
    ("Five interpolation engines", "Custom UNetVFI + RIFE + FILM + Super-SloMo + RAFT, with a classical "
     "TV-L1 baseline for honest comparison.", BLUE),
    ("Temporal upscaling", "Recursive ×2/×4 (30→15→7.5) plus Continuous mode — any cadence, each frame "
     "direct from two real frames (no error compounding).", TEAL),
    ("Standards I/O", "Reads .nc/.h5 (GOES / Himawari / INSAT), writes CF NetCDF brightness "
     "temperature — the PS contract.", VIOLET),
    ("Validated vs ground truth", "PSNR · SSIM · FSIM · MSE · MAE(K) · LPIPS + cloud-motion metrics "
     "(flow-EPE, edge-SSIM, temporal warping).", ORANGE),
    ("Web dashboard", "Three tabs — Interpolate · Temporal Upscaling · Validation Report — with "
     "animations, motion overlay, live metrics.", GREEN),
    ("Cloud-ready", "One command connects to a free T4 (Lightning.ai / Colab / Kaggle) with 100 GB "
     "persistent storage.", NAVY),
]
xs = [0.55, 3.72, 6.89]; ys = [1.85, 3.32]; w, h = 2.95, 1.32
for k, (t_, b_, col) in enumerate(feats):
    x = xs[k % 3]; y = ys[k // 3]
    rect(sl, x, y, w, h, WHITE, line=HAIR, radius=0.08)
    badge(sl, x + 0.32, y + 0.32, k + 1, col, d=0.34)
    tb(sl, x + 0.58, y + 0.12, w - 0.7, 0.4, [{"t": t_, "s": 12.5, "c": INK, "b": True}])
    tb(sl, x + 0.18, y + 0.55, w - 0.34, 0.72, [{"t": b_, "s": 10.3, "c": SLATE}])
rect(sl, 0.55, 4.85, 9.0, 0.42, NAVY, radius=0.2)
settext(sl.shapes[-1], [{"runs": [
    {"t": "Metric  ", "s": 11, "c": ORANGE, "b": True},
    {"t": "PSNR = 10·log₁₀(1 / MSE)  →  perfect frame → PSNR infinite.   Our UNetVFI reached val PSNR ≈ 42 dB, "
          "SSIM ≈ 0.95 on held-out GOES.", "s": 11, "c": WHITE}]}], anchor=MSO_ANCHOR.MIDDLE)

# ===============================================================================
# SLIDE 5 — Process flow (clean stages + thin arrowed connectors + equation)
sl = S[4]; clear_prompts(sl)
head(sl, "03 · How it works", "Process flow", accent=TEAL)
stages = [
    (BLUE, "INPUT", ["Two real frames", "I(t₀), I(t₂)", ".nc / .h5 · TIR ~10µm"]),
    (TEAL, "PREP", ["Calibrate → BT (K)", "Normalize → [0,1]", "Tile 256²"]),
    (VIOLET, "AI OPTICAL FLOW", ["Estimate motion", "+ synthesise I(t)", "t = ½, ¼, ¾"]),
    (GREEN, "REBUILD", ["Untile (feather)", "Denormalize → BT", "Write .nc @ new t"]),
    (ORANGE, "VALIDATE & SHOW", ["PSNR/SSIM/FSIM", "web dashboard", "time-lapse"]),
]
w, gap, y, h = 1.55, 0.31, 2.05, 1.75
x = 0.55
for i, (col, hd, lines) in enumerate(stages):
    rect(sl, x, y, w, h, WHITE, line=HAIR, radius=0.09)
    rect(sl, x, y, w, 0.07, col, rounded=False)
    badge(sl, x + 0.3, y + 0.4, i + 1, col, d=0.3)
    tb(sl, x + 0.5, y + 0.22, w - 0.55, 0.34, [{"t": hd, "s": 10.5, "c": col, "b": True}])
    tb(sl, x + 0.15, y + 0.7, w - 0.3, 1.0,
       [{"t": ln, "s": 9.8, "c": SLATE, "sa": 2} for ln in lines])
    if i < len(stages) - 1:
        conn(sl, x + w + 0.04, y + h / 2, x + w + gap - 0.04, y + h / 2, INK, w=1.4)
    x += w + gap
rect(sl, 0.55, 4.1, 9.0, 1.05, MIST, line=HAIR, radius=0.06)
tb(sl, 0.8, 4.18, 8.6, 0.9, [
    {"t": "Synthesis equation (per pixel x)", "s": 10.5, "c": TEAL, "b": True, "sa": 3},
    {"t": "I_t(x) ≈ M(x) · I₀(x + F_{t→0}(x))   +   (1 − M(x)) · I₂(x + F_{t→1}(x))", "s": 15, "c": INK,
     "b": True, "sa": 3},
    {"t": "F = learned optical flow (motion vectors) · M = learned visibility / occlusion mask · "
          "t = (t₁−t₀)/(t₂−t₀);  for 30→15 min,  t = 0.5", "s": 10, "c": SLATE, "i": True}])

# ===============================================================================
# SLIDE 6 — realistic dashboard mock
sl = S[5]; clear_prompts(sl)
head(sl, "04 · Visualisation", "Web dashboard (GUI)", accent=ORANGE)


def irframe(sl, l, t, w, h, accent=None, tag=None):
    rect(sl, l, t, w, h, IRBG, line=(accent or RGBColor(0x39, 0x44, 0x55)), lw=(1.75 if accent else 0.75), radius=0.05)
    blobs = [(0.12, 0.30, 0.42, 0.34, RGBColor(0xDD, 0xE4, 0xEC)),
             (0.40, 0.12, 0.34, 0.30, RGBColor(0xAF, 0xBC, 0xCB)),
             (0.55, 0.50, 0.40, 0.40, RGBColor(0xC7, 0xD1, 0xDC)),
             (0.20, 0.62, 0.26, 0.26, RGBColor(0x93, 0xA2, 0xB4))]
    for bx, by, bw, bh, col in blobs:
        oval(sl, l + bx * w, t + by * h, bw * w, bh * h, col)
    if tag:
        rect(sl, l + 0.04, t + 0.04, 0.5, 0.2, (accent or RGBColor(0x39, 0x44, 0x55)), rounded=False)
        settext(sl.shapes[-1], [{"t": tag, "s": 8, "c": WHITE, "b": True, "align": PP_ALIGN.CENTER}],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# window
rect(sl, 0.55, 1.8, 9.0, 3.35, WHITE, line=RGBColor(0xC9, 0xD3, 0xDE), lw=1.0, radius=0.03)
rect(sl, 0.55, 1.8, 9.0, 0.34, NAVY, radius=0.03)
for i, c in enumerate([RGBColor(0xF2, 0x6A, 0x1B), GREEN, BLUE]):
    oval(sl, 0.72 + i * 0.17, 1.9, 0.11, 0.11, c)
tb(sl, 1.4, 1.83, 7.0, 0.28, [{"t": "PS-12 · Satellite Frame Interpolation          localhost:8501",
                               "s": 9.5, "c": WHITE, "b": True}], anchor=MSO_ANCHOR.MIDDLE)
# sidebar
rect(sl, 0.72, 2.3, 1.95, 2.7, MIST, line=HAIR, radius=0.05)
tb(sl, 0.85, 2.4, 1.75, 2.5, [
    {"t": "CONFIGURATION", "s": 9, "c": SLATE, "b": True, "sa": 6},
    {"t": "Source:    GOES-19", "s": 9.5, "c": INK, "sa": 6},
    {"t": "Model:     UNetVFI", "s": 9.5, "c": INK, "sa": 6},
    {"t": "Factor:    2× | 4×", "s": 9.5, "c": INK, "sa": 6},
    {"t": "Motion overlay:  ON", "s": 9.5, "c": INK, "sa": 8},
    {"t": "RUNNABLE", "s": 9, "c": SLATE, "b": True, "sa": 4},
    {"t": "● unet   ● raft   ● classical", "s": 9, "c": GREEN, "b": True}])
# tabs
for i, (lab, c, on) in enumerate([("Interpolate", BLUE, True), ("Temporal Upscaling", TEAL, False),
                                  ("Validation Report", ORANGE, False)]):
    rect(sl, 2.85 + i * 1.9, 2.34, 1.78, 0.32, c if on else WHITE, line=(None if on else HAIR), radius=0.2)
    settext(sl.shapes[-1], [{"t": lab, "s": 9.5, "c": (WHITE if on else SLATE), "b": True,
            "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# frame strip
fr = [("t₀  input", None), ("t = 0.5  AI", BLUE), ("t₂  input", None)]
for i, (lab, ac) in enumerate(fr):
    fx = 2.9 + i * 2.18
    irframe(sl, fx, 2.78, 2.0, 1.32, accent=ac, tag=("AI" if ac else None))
    tb(sl, fx, 4.12, 2.0, 0.24, [{"t": lab, "s": 9, "c": (BLUE if ac else SLATE), "b": bool(ac),
                                  "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    if i < 2:
        conn(sl, fx + 2.0, 3.44, fx + 2.18, 3.44, MUTE, w=1.2)
# metric chips
for i, (lab, val, c) in enumerate([("PSNR", "33.5", BLUE), ("SSIM", "0.89", TEAL), ("FSIM", "0.99", VIOLET),
                                   ("MAE(K)", "1.7", ORANGE)]):
    cx = 2.9 + i * 1.62
    rect(sl, cx, 4.45, 1.5, 0.5, WHITE, line=HAIR, radius=0.12)
    settext(sl.shapes[-1], [{"t": val, "s": 13, "c": c, "b": True, "sa": 0, "align": PP_ALIGN.CENTER},
            {"t": lab, "s": 8, "c": SLATE, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
tb(sl, 0.55, 5.18, 9.0, 0.3, [{"t": "Served from the cloud T4 → opened in your browser (Lightning "
        "Streamlit plugin / ngrok). Validation tab renders committed comparison results.", "s": 9.5,
        "c": MUTE, "i": True}])

# ===============================================================================
# SLIDE 7 — Architecture (research-paper U-Net)
sl = S[6]; clear_prompts(sl)
head(sl, "05 · Architecture", "Custom UNetVFI — flow + visibility", accent=VIOLET)


def fmap(sl, cx, cyc, h, ch, fill, res):
    w = 0.5
    rect(sl, cx - w / 2, cyc - h / 2, w, h, fill, line=RGBColor(0x2A, 0x3A, 0x55), lw=0.5, rounded=False)
    tb(sl, cx - 0.45, cyc + h / 2 + 0.02, 0.9, 0.2, [{"t": ch, "s": 8.5, "c": INK, "b": True,
       "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    tb(sl, cx - 0.45, cyc - h / 2 - 0.22, 0.9, 0.18, [{"t": res, "s": 7.5, "c": MUTE,
       "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)


cy = 3.0
# inputs + concat
irframe(sl, 0.55, 2.45, 0.62, 0.5)
tb(sl, 0.55, 2.96, 0.62, 0.18, [{"t": "I₀", "s": 9, "c": INK, "b": True, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
irframe(sl, 0.55, 3.15, 0.62, 0.5)
tb(sl, 0.55, 3.66, 0.62, 0.18, [{"t": "I₂", "s": 9, "c": INK, "b": True, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
# t-plane: the target time t fed in as a 3rd input channel (time-conditioned / arbitrary-time)
rect(sl, 0.55, 3.9, 0.62, 0.32, TEAL, radius=0.1)
settext(sl.shapes[-1], [{"t": "t-plane", "s": 7.5, "c": WHITE, "b": True, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(sl, 1.3, 2.72, 0.58, 0.66, NAVY, radius=0.12)
settext(sl.shapes[-1], [{"t": "concat", "s": 8.5, "c": WHITE, "b": True, "sa": 0, "align": PP_ALIGN.CENTER},
        {"t": "3 ch", "s": 7, "c": WHITE, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
conn(sl, 1.17, 2.7, 1.34, 2.9, MUTE, w=1.0); conn(sl, 1.17, 3.4, 1.34, 3.1, MUTE, w=1.0)
conn(sl, 1.17, 4.05, 1.34, 3.3, TEAL, w=1.0, dash="dash")   # t-plane -> concat
# encoder (down-step) E1..E4
encx = [2.15, 2.9, 3.65, 4.4]; ench = [1.5, 1.15, 0.85, 0.6]
encmeta = [("32", "256²"), ("64", "128²"), ("128", "64²"), ("256", "32²")]
for i, ex in enumerate(encx):
    fmap(sl, ex, cy, ench[i], encmeta[i][0], ENC[i], encmeta[i][1])
conn(sl, 1.85, 3.05, encx[0] - 0.27, 3.0, INK, w=1.3)
for i in range(3):
    conn(sl, encx[i] + 0.27, cy, encx[i + 1] - 0.27, cy, INK, w=1.3)
# decoder (up-step) D3..D1
decx = [5.25, 6.0, 6.75]; dech = [0.85, 1.15, 1.5]
decmeta = [("128", "64²"), ("64", "128²"), ("32", "256²")]
for i, dx in enumerate(decx):
    fmap(sl, dx, cy, dech[i], decmeta[i][0], DEC[i], decmeta[i][1])
conn(sl, encx[3] + 0.27, cy, decx[0] - 0.27, cy, INK, w=1.3)
for i in range(2):
    conn(sl, decx[i] + 0.27, cy, decx[i + 1] - 0.27, cy, INK, w=1.3)
# skip connections (staggered staples over the top, dashed)
skips = [(encx[0], decx[2], 1.78), (encx[1], decx[1], 1.92), (encx[2], decx[0], 2.06)]
for ex, dx, ytop in skips:
    conn(sl, ex, cy - ench[encx.index(ex)] / 2, ex, ytop, MUTE, w=0.9, head=False, dash="dash")
    conn(sl, ex, ytop, dx, ytop, MUTE, w=0.9, head=False, dash="dash")
    conn(sl, dx, ytop, dx, cy - dech[decx.index(dx)] / 2, MUTE, w=0.9, dash="dash")
tb(sl, 4.3, 1.6, 2.0, 0.2, [{"t": "skip connections", "s": 8.5, "c": MUTE, "i": True, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
# head -> flow + mask + source(PINN)
conn(sl, decx[2] + 0.27, cy, 7.22, cy, INK, w=1.3)
heads = [("flow t→0", ORANGE, "2ch"), ("flow t→1", ORANGE, "2ch"),
         ("mask M", RGBColor(0xD4, 0xA0, 0x17), "1ch · σ"), ("source S", TEAL, "1ch · PINN")]
for i, (lab, c, sub) in enumerate(heads):
    rect(sl, 7.28, 1.92 + i * 0.5, 1.22, 0.44, c, radius=0.16)
    settext(sl.shapes[-1], [{"t": lab, "s": 9, "c": WHITE, "b": True, "sa": 0, "align": PP_ALIGN.CENTER},
            {"t": sub, "s": 7, "c": WHITE, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
conn(sl, 8.5, 2.62, 8.72, 2.62, INK, w=1.3)                       # flows+mask -> warp
conn(sl, 8.5, 3.42, 8.72, 3.42, TEAL, w=1.1, dash="dash")          # source -> PINN loss (training only)
rect(sl, 8.74, 2.25, 0.92, 0.78, NAVY, radius=0.1)
settext(sl.shapes[-1], [{"t": "warp", "s": 9, "c": WHITE, "b": True, "sa": 0, "align": PP_ALIGN.CENTER},
        {"t": "+ blend", "s": 8.5, "c": WHITE, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
irframe(sl, 8.9, 3.2, 0.6, 0.46, accent=GREEN)
tb(sl, 8.72, 3.66, 0.95, 0.2, [{"t": "I(t)", "s": 9.5, "c": GREEN, "b": True, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
conn(sl, 9.2, 3.05, 9.2, 3.18, GREEN, w=1.3)
tb(sl, 8.55, 3.86, 1.1, 0.2, [{"t": "PINN loss", "s": 7.5, "c": TEAL, "i": True, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
# caption strip
rect(sl, 0.55, 4.5, 9.0, 0.9, MIST, line=HAIR, radius=0.06)
tb(sl, 0.75, 4.55, 8.7, 0.85, [
    {"runs": [{"t": "Blend:  ", "s": 9.5, "c": VIOLET, "b": True},
              {"t": "I(t) = M·warp(I₀, F_{t→0}) + (1−M)·warp(I₂, F_{t→1})", "s": 9.5, "c": INK, "b": True},
              {"t": "   — RIFE-style intermediate flow + Super-SloMo visibility mask.", "s": 9.5, "c": SLATE}], "sa": 2},
    {"runs": [{"t": "Physics-informed (PINN, training):  ", "s": 9.5, "c": TEAL, "b": True},
              {"t": "advection ∂I/∂t + u·grad(I) = S — source S models cloud growth; +0.12 PSNR measured.",
               "s": 9.5, "c": SLATE}], "sa": 2},
    {"runs": [{"t": "Arbitrary-time (time-conditioned):  ", "s": 9.5, "c": TEAL, "b": True},
              {"t": "target time t is an input channel; trained on a configurable t-grid × gap-granule "
                    "set → one model renders any time 30→15→7.5 min (t=0.25/0.5/0.75), not just the midpoint.",
               "s": 9.5, "c": SLATE}], "sa": 2},
    {"runs": [{"t": "~2–5M params · trains on a T4 in hours · self-supervised on INSAT · "
                    "Local → GitHub → Lightning T4 → Streamlit.", "s": 9, "c": SLATE}]}])

# ===============================================================================
# SLIDE 8 — Technologies (clean rows, colored dot, no bars)
sl = S[7]; clear_prompts(sl)
head(sl, "06 · Stack", "Technologies used", accent=BLUE)
groups = [
    ("Deep learning", BLUE, "PyTorch · RAFT (torchvision) · RIFE · FILM · Super-SloMo · custom UNetVFI"),
    ("Satellite I/O", TEAL, "xarray · netCDF4 · h5py · satpy · boto3 (GOES S3) · paramiko / lftp (MOSDAC SFTP)"),
    ("Classical + metrics", VIOLET, "OpenCV TV-L1 / Farnebäck · scikit-image · piq (FSIM / LPIPS)"),
    ("Web + serving", ORANGE, "Streamlit (3-tab dashboard) · ngrok / Lightning Streamlit plugin"),
    ("Compute + DevOps", GREEN, "Lightning.ai / Colab / Kaggle free T4 + 100 GB persist · Git / GitHub"),
    ("Datasets", NAVY, "GOES-19 ABI Ch13 (10.3µm) · Himawari AHI B13 (10.4µm) · INSAT-3DR/3DS TIR1 (10.8µm)"),
]
y = 1.85
for name, col, items in groups:
    oval(sl, 0.6, y + 0.13, 0.16, 0.16, col)
    tb(sl, 0.9, y, 8.6, 0.46, [{"runs": [{"t": name + "   ", "s": 13, "c": col, "b": True},
                                         {"t": items, "s": 12, "c": INK}]}])
    if y < 4.5:
        rect(sl, 0.9, y + 0.5, 8.6, 0.012, HAIR, rounded=False)
    y += 0.56

# ===============================================================================
# SLIDE 9 — Cost
sl = S[8]; clear_prompts(sl)
head(sl, "07 · Feasibility", "Estimated cost", accent=GREEN)
rect(sl, 0.55, 1.95, 3.0, 2.9, NAVY, radius=0.06)
settext(sl.shapes[-1], [
    {"t": "₹0", "s": 46, "c": WHITE, "b": True, "sa": 2, "align": PP_ALIGN.CENTER},
    {"t": "on free tiers", "s": 13, "c": RGBColor(0xBE, 0xD3, 0xE6), "sa": 8, "align": PP_ALIGN.CENTER},
    {"t": "≤ $10 worst case", "s": 12, "c": ORANGE, "b": True, "align": PP_ALIGN.CENTER}],
    anchor=MSO_ANCHOR.MIDDLE)
items = [("Software", "Fully open-source.", BLUE),
         ("Data", "NOAA GOES / Himawari on AWS (free); INSAT via MOSDAC (free account).", TEAL),
         ("Compute", "Free T4 — Lightning.ai (~35 hrs/mo), Colab, Kaggle.", GREEN),
         ("If paid", "Training ≈ 2–5 T4-GPU-hours.", VIOLET),
         ("Optional", "OpenAI report narrative < $10 (not required).", ORANGE)]
y = 1.95
for n, d, c in items:
    rect(sl, 3.85, y, 5.7, 0.5, MIST, line=HAIR, radius=0.12)
    settext(sl.shapes[-1], [{"runs": [{"t": n + ":  ", "s": 11.5, "c": c, "b": True},
            {"t": d, "s": 11.5, "c": INK}]}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.58
tb(sl, 3.85, 4.85, 5.7, 0.3, [{"runs": [{"t": "Cost model:  ", "s": 11, "c": GREEN, "b": True},
        {"t": "cost ≈ GPU-hrs × $/hr ≈ 5 × $0.35 ≈ $1.75 per full training run.", "s": 11, "c": SLATE, "i": True}]}])

# ===============================================================================
# SLIDE 10 — Outcomes (NO "thank you")
sl = S[9]
tb(sl, 0.52, 0.72, 9.0, 0.26, [{"t": "OUTCOMES", "s": 10.5, "c": ORANGE, "b": True}])
tb(sl, 0.5, 0.98, 9.0, 0.56, [{"t": "Validated, reproducible, and ready", "s": 23, "c": INK, "b": True}])
rect(sl, 0.54, 1.55, 1.25, 0.05, ORANGE, rounded=False)
outs = [
    ("Real GPU run", "Trained on a Tesla T4; deterministic battery 16/16 green.", BLUE),
    ("Measured quality", "UNetVFI val PSNR ≈ 42 / SSIM ≈ 0.95; full classical / RAFT / UNet comparison.", TEAL),
    ("End-to-end", ".nc → .nc interpolation + 30→15→7.5 min upscaling + web dashboard.", VIOLET),
    ("Open & free", "Open-source backbones, free-tier T4, datasets at ₹0.", GREEN),
]
for i, (t_, d_, c) in enumerate(outs):
    x = 0.55 + (i % 2) * 4.62; y = 1.85 + (i // 2) * 1.35
    rect(sl, x, y, 4.42, 1.18, WHITE, line=HAIR, radius=0.08)
    rect(sl, x, y, 0.09, 1.18, c, rounded=False) if False else None  # (intentionally no left bar)
    oval(sl, x + 0.22, y + 0.22, 0.18, 0.18, c)
    tb(sl, x + 0.55, y + 0.14, 3.7, 0.4, [{"t": t_, "s": 13.5, "c": INK, "b": True}])
    tb(sl, x + 0.25, y + 0.58, 4.0, 0.55, [{"t": d_, "s": 11, "c": SLATE}])
rect(sl, 0.55, 4.75, 9.0, 0.5, NAVY, radius=0.18)
settext(sl.shapes[-1], [{"runs": [
    {"t": "Repository   ", "s": 12, "c": ORANGE, "b": True},
    {"t": "github.com/uditsenapaty/ps12", "s": 13, "c": WHITE, "b": True},
    {"t": "      Enhancing temporal resolution — no new satellite required.", "s": 11, "c": RGBColor(0xC6, 0xD6, 0xE8), "i": True}]}],
    anchor=MSO_ANCHOR.MIDDLE)

OUT = r"D:\Udit\gitclones\ps12\ISRO_BAH_2026_PS12 - FILLED.pptx"
try:
    prs.save(OUT)
except PermissionError:
    OUT = OUT.replace(".pptx", " (rebuilt).pptx")  # main file is open in PowerPoint
    prs.save(OUT)
    print("NOTE: main file was locked (open in PowerPoint) — saved to a new file.")
print("SAVED:", OUT, "| slides:", len(S))
