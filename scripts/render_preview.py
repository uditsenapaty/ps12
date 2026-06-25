#!/usr/bin/env python
"""Render the filled deck to PNGs for visual QA (no LibreOffice needed).

Pastes the real template background and draws every shape (rect/oval/connector/text) at its exact
coordinates with Pillow. Not pixel-perfect (approx fonts/wrap), but enough to judge layout + balance.
"""
import glob, os
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

W, H = 1600, 900
PXIN = 160.0  # px per inch
MEDIA = r"D:\Udit\gitclones\ps12\_tpl_media"
OUT = r"D:\Udit\gitclones\ps12\_preview"; os.makedirs(OUT, exist_ok=True)
FILE = [s for s in glob.glob(r"D:\Udit\gitclones\ps12\*.pptx") if "FILLED" in s][0]


def font(sz, bold=False):
    for nm in ([r"C:\Windows\Fonts\segoeuib.ttf"] if bold else [r"C:\Windows\Fonts\segoeui.ttf"]) + [r"C:\Windows\Fonts\arial.ttf"]:
        try:
            return ImageFont.truetype(nm, int(sz))
        except Exception:
            continue
    return ImageFont.load_default()


def emu_px(v):
    return None if v is None else Emu(v).inches * PXIN


def rgb(color):
    try:
        return tuple(bytes.fromhex(str(color.rgb)))
    except Exception:
        return None


def wrap(draw, text, fnt, maxw):
    out, cur = [], ""
    for word in text.split():
        test = (cur + " " + word).strip()
        if draw.textlength(test, font=fnt) <= maxw or not cur:
            cur = test
        else:
            out.append(cur); cur = word
    if cur:
        out.append(cur)
    return out


def bg_for(i):
    m = {1: "image2.png", 2: "image3.png"}.get(i, "image1.png")
    p = os.path.join(MEDIA, m)
    base = Image.new("RGB", (W, H), "white")
    if os.path.exists(p):
        im = Image.open(p).convert("RGBA").resize((W, H))
        base.paste(im, (0, 0), im)  # composite over white = how PowerPoint shows the transparent body
    return base


prs = Presentation(FILE)
for i, sl in enumerate(prs.slides, 1):
    img = bg_for(i); d = ImageDraw.Draw(img, "RGBA")
    for sh in sl.shapes:
        l, t, w, h = emu_px(sh.left), emu_px(sh.top), emu_px(sh.width), emu_px(sh.height)
        if None in (l, t, w, h):
            continue
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.LINE or sh.__class__.__name__ == "Connector":
            col = rgb(sh.line.color) or (90, 100, 120)
            x1, y1, x2, y2 = l, t, l + w, t + h
            try:
                fh = sh._element.spPr.xfrm.get("flipH") == "1"; fv = sh._element.spPr.xfrm.get("flipV") == "1"
                if fh:
                    x1, x2 = x2, x1
                if fv:
                    y1, y2 = y2, y1
            except Exception:
                pass
            d.line([x1, y1, x2, y2], fill=col, width=2)
            continue
        fill = None
        try:
            if sh.fill.type is not None:
                fill = rgb(sh.fill.fore_color)
        except Exception:
            fill = None
        line = None
        try:
            line = rgb(sh.line.color)
        except Exception:
            line = None
        is_oval = (st == MSO_SHAPE_TYPE.AUTO_SHAPE and getattr(sh, "auto_shape_type", None) is not None
                   and str(sh.auto_shape_type) == "OVAL")
        if fill or line:
            box = [l, t, l + w, t + h]
            if is_oval:
                d.ellipse(box, fill=fill, outline=line, width=1)
            else:
                try:
                    d.rounded_rectangle(box, radius=min(12, h / 3), fill=fill, outline=line, width=1)
                except Exception:
                    d.rectangle(box, fill=fill, outline=line, width=1)
        # text
        if sh.has_text_frame and sh.text_frame.text.strip():
            tf = sh.text_frame; cy = t + 5
            for para in tf.paragraphs:
                txt = "".join(r.text for r in para.runs) or para.text
                if not txt.strip():
                    cy += 8; continue
                run = para.runs[0] if para.runs else None
                sz = (run.font.size.pt if run and run.font.size else 12) * PXIN / 72.0
                bold = bool(run.font.bold) if run else False
                tcol = rgb(run.font.color) if run else None
                tcol = tcol or (20, 30, 55)
                fnt = font(sz, bold)
                align = str(para.alignment) if para.alignment else "LEFT"
                for ln in wrap(d, txt, fnt, w - 10):
                    tw = d.textlength(ln, font=fnt)
                    x = l + 6
                    if "CENTER" in align:
                        x = l + (w - tw) / 2
                    elif "RIGHT" in align:
                        x = l + w - tw - 6
                    d.text((x, cy), ln, font=fnt, fill=tcol)
                    cy += sz * 1.2
    img.save(os.path.join(OUT, f"slide{i:02d}.png"))
print("rendered", len(prs.slides.__iter__.__self__._sldIdLst), "-> ", OUT)
